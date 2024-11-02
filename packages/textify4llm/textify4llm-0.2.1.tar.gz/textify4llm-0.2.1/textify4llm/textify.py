import os
import json
from pptx import Presentation
from docx import Document
import pandas as pd
import pdfplumber
import fitz
from PIL import Image
import io
import torch
from paddleocr import PaddleOCR
from transformers import BlipProcessor, BlipForConditionalGeneration
import logging
import numpy as np
import re

os.environ["TOKENIZERS_PARALLELISM"] = "false" # supress warning from BLIP
logging.getLogger("ppocr").setLevel(logging.ERROR) # supress paddle log

def process_file(file_path):

    xml_based_files = [
        'docx',     # Microsoft Word
        'xlsx',     # Microsoft Excel
        'pptx',     # Microsoft PowerPoint
        'pdf'   
    ]

    image_files = [
        'jpg',      # JPEG images
        'png',      # PNG images
        'jpeg',     # Same as JPG
        'bmp',      # Bitmap images
    ]

    file_extension = os.path.splitext(file_path)[-1].lower().strip('.')

    if file_extension in xml_based_files:
        # 这几个XML还都不一样
        # import各自的库
        return _handle_xml(file_path, file_extension)
    elif file_extension in image_files:
        return _handle_image(file_path=file_path)
    else:
        # 把其他格式一揽子按照纯文本试一下
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return _handle_text(file_path)
        except (UnicodeDecodeError, FileNotFoundError, IsADirectoryError):
            return 'Unsupported file type'
        
def _handle_text(file_path):
    '''Handles text-based files.'''
    with open(file_path, 'r', encoding='utf-8') as file:
        content = file.read()
    file_name = os.path.basename(file_path)
    text_data = {
        'filename': file_name,
        'content': content
    }
    return text_data

def _handle_xml(file_path, file_extension):
    '''Handles XML-based files.'''
    if file_extension == 'docx':
        return _handle_docx(file_path)
    elif file_extension == 'xlsx':
        return _handle_xlsx(file_path)
    elif file_extension == 'pptx':
        return _handle_pptx(file_path)
    elif file_extension == 'pdf':
        return _handle_pdf(file_path)
    
def _handle_docx(file_path):
    '''Handles DOCX files.'''
    file_name = os.path.basename(file_path)
    doc_data = {'filename': file_name, 'content': []}

    try:
        doc = Document(file_path)
        sections = doc.sections
        page_size = None
        margins = None

        # Extract page size and margins if available
        if sections:
            first_section = sections[0]
            page_size = {
                'width': first_section.page_width.pt if first_section.page_width else None,
                'height': first_section.page_height.pt if first_section.page_height else None
            }
            margins = {
                'top_margin': first_section.top_margin.pt if first_section.top_margin else None,
                'bottom_margin': first_section.bottom_margin.pt if first_section.bottom_margin else None,
                'left_margin': first_section.left_margin.pt if first_section.left_margin else None,
                'right_margin': first_section.right_margin.pt if first_section.right_margin else None
            }

        # Extract paragraphs with formatting details and embedded images
        for para_number, para in enumerate(doc.paragraphs):
            para_data = {
                'paragraph': para_number + 1,
                'text': para.text.strip(),
                'font': [run.font.name for run in para.runs if run.font],
                'font_size': [run.font.size.pt if run.font.size else None for run in para.runs],
                'bold': [run.bold for run in para.runs],
                'italic': [run.italic for run in para.runs],
                'underline': [run.underline for run in para.runs],
                'line_spacing': para.paragraph_format.line_spacing if para.paragraph_format.line_spacing else None
            }
            doc_data['content'].append(para_data)

        # Extract images from the document
        images = []
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                image_data = rel.target_part.blob
                images.append(_handle_image(blob=image_data))

        # Include page size, margins, and images information
        if page_size:
            doc_data['page_size'] = page_size
        if margins:
            doc_data['margins'] = margins
        if images:
            doc_data['images'] = images

        return doc_data
    except Exception as e:
        return {'error': f'Error processing DOCX: {str(e)}'}


def _handle_xlsx(file_path):
    '''Handles XLSX files.'''
    try:
        file_name = os.path.basename(file_path)
        xlsx_data = {'filename': file_name, 'content': []}
        
        # Load the Excel file
        xls = pd.ExcelFile(file_path)
        for sheet_name in xls.sheet_names:
            sheet = xls.parse(sheet_name)
            csv_data = sheet.to_csv(sep='\t', index=False)  # Convert each sheet to TSV format
            xlsx_data['content'].append({
                'sheet_name': sheet_name,
                'data': csv_data
            })
        
        return xlsx_data
    except Exception as e:
        return {'error': f'Error processing XLSX: {str(e)}'}
    
def _handle_pptx(file_path):
    '''Handles PPT files.'''
    try:
        file_name = os.path.basename(file_path)
        ppt_data = {'filename': file_name, 'content': []}

        presentation = Presentation(file_path)
        
        for slide_number, slide in enumerate(presentation.slides):
            slide_text = " ".join([shape.text for shape in slide.shapes if hasattr(shape, "text")]).strip()
            slide_images = [_handle_image(blob=shape.image.blob) for shape in slide.shapes if shape.shape_type == 13]  # Shape type 13 is Picture
            
            ppt_data['content'].append({
                'slide': slide_number + 1,
                'text': slide_text,
                'images': slide_images
            })

        return ppt_data
    except Exception as e:
        return {'error': f'Error processing PPTX: {str(e)}'}


def _handle_pdf(file_path):
    '''Handles PDF files.'''
    file_name = os.path.basename(file_path)
    pdf_data = {'filename': file_name, 'content': []}
    try: 
        with pdfplumber.open(file_path) as pdf, fitz.open(file_path) as doc:
                for page_number, page in enumerate(pdf.pages):
                    text = page.extract_text().strip() if page.extract_text() else ""
                    tables = page.extract_tables()
                    images = []
                    img_list = doc[page_number].get_images(full=True)
                    for img in img_list:
                        xref = img[0]
                        base_image = doc.extract_image(xref)
                        images.append(_handle_image(blob=base_image['image']))
                    pdf_data['content'].append({
                        'page': page_number + 1,
                        'text': text,
                        'images': images,
                        'tables': tables
                    })
    
        return pdf_data
    except Exception as e:
        return {'error': f'Error processing PDF: {str(e)}'}

def _handle_image(file_path=None, xref=None, blob=None):
    '''Handles image-based files, performs OCR in both Chinese and English, and generates captions.'''
    
    image = _build_image(file_path=file_path, xref=xref, blob=blob)
    if image is None:
        return {"error": "Invalid image input"}
    
    ocr = PaddleOCR(lang='ch', use_angle_cls=False)
    image_np = np.array(image)
    result = ocr.ocr(image_np, cls=True)  # Set `cls=True` for better accuracy in Chinese text recognition
    if result[0]:
        extracted_text = " ".join([line[1][0] for line in result[0]]).strip()
    else:
        extracted_text = ""

    blip_processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base")
    blip_model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-base")
    inputs = blip_processor(image, return_tensors="pt")
    
    with torch.no_grad():
        out = blip_model.generate(**inputs, max_new_tokens=50)
    
    caption = blip_processor.decode(out[0], skip_special_tokens=True)

    return {
        'content': extracted_text,
        'caption': caption
    }


def _build_image(file_path=None, xref=None, blob=None):
    if file_path:
        return Image.open(file_path)
    elif xref is not None:
        return Image.open(io.BytesIO(xref))
    elif blob is not None:
        return Image.open(io.BytesIO(blob))
    return None